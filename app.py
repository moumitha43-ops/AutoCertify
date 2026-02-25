import streamlit as st
import pandas as pd
from pptx import Presentation
from pathlib import Path
import comtypes.client
import smtplib
from email.message import EmailMessage
from email.utils import make_msgid
import tempfile
import pythoncom

pythoncom.CoInitialize()

def fill_ppt(template_path, output_ppt, data):
    prs = Presentation(str(template_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in data.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(
                                placeholder,
                                "" if pd.isna(value) else str(value)
                            )

    prs.save(str(output_ppt))


def ppt_to_images(ppt_path, output_folder):
    ppt_path = Path(ppt_path).resolve()
    output_folder = Path(output_folder).resolve()
    output_folder.mkdir(parents=True, exist_ok=True)

    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT not found: {ppt_path}")

    pythoncom.CoInitialize()

    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(str(ppt_path))
        presentation.SaveAs(str(output_folder), 18)  # 18 = PNG
        presentation.Close()
        powerpoint.Quit()

    finally:
        pythoncom.CoUninitialize()


def send_email(sender, password, recipient, subject, body, attachments, name, event_name):
    msg = EmailMessage()
    msg["From"] = sender
    msg["To"] = recipient
    msg["Subject"] = subject
    msg.set_content(body)

    cid = make_msgid()

    html_content = f"""
    <html>
      <body style="font-family: Arial">
        <p>Dear <b>{name}</b>,</p>
        <p>Please find your certificate below:</p>
        <img src="cid:{cid[1:-1]}" width="700">
        <p>Regards,<br>Event Team</p>
      </body>
    </html>
    """

    msg.add_alternative(html_content, subtype="html")

    # Attach first image inline
    if attachments:
        with open(attachments[0], "rb") as img:
            img_data = img.read()

        msg.get_payload()[1].add_related(
            img_data,
            maintype="image",
            subtype="png",
            cid=cid,
            filename=f"{event_name}.png",
            disposition="inline"
        )

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(sender, password)
        server.send_message(msg)


st.set_page_config(page_title="Certificate Generator", layout="centered")
st.title("🎓 Event Certificate Generator")

event_name = st.text_input("📌 Event Name")

ppt_file = st.file_uploader("📄 Upload PowerPoint Template", type="pptx")
csv_file = st.file_uploader("📊 Upload CSV File", type="csv")

st.subheader("📧 Email Configuration")
sender_email = st.text_input("Sender Email")
sender_password = st.text_input("Sender App Password", type="password")
email_column = st.text_input("CSV Column Name for Recipient Email", value="EMAIL")

if csv_file is not None:
    csv_file.seek(0)
    df = pd.read_csv(csv_file)

    if "NAME" in df.columns:
        selected_name = st.selectbox(
            "Select a name to preview",
            df["NAME"].astype(str).tolist()
        )

        if st.button("👁 Generate Preview") and ppt_file is not None:
            row = df[df["NAME"].astype(str) == selected_name].iloc[0].to_dict()

            with tempfile.TemporaryDirectory() as tmp:
                template = Path(tmp) / "template.pptx"
                template.write_bytes(ppt_file.read())

                ppt_out = Path(tmp) / "preview.pptx"
                img_out = Path(tmp) / "images"

                fill_ppt(template, ppt_out, row)
                ppt_to_images(ppt_out, img_out)

                images = list(img_out.glob("*.PNG"))
                if images:
                    st.image(str(images[0]), caption=f"Preview for {selected_name}", use_column_width=True)


if st.button("🚀 Generate & Send Certificates"):

    if not all([event_name, ppt_file, csv_file, sender_email, sender_password]):
        st.error("❌ Please fill all fields")
        st.stop()
    csv_file.seek(0)
    df = pd.read_csv(csv_file)

    total = len(df)
    progress_bar = st.progress(0)
    status_text = st.empty()

    base_output = Path("output") / event_name.replace(" ", "_")
    base_output.mkdir(parents=True, exist_ok=True)

    sent_count = 0

    with tempfile.TemporaryDirectory() as tmpdir:
        template_path = Path(tmpdir) / "template.pptx"
        template_path.write_bytes(ppt_file.read())

        for idx, row in df.iterrows():
            data = row.to_dict()
            name = str(data.get("NAME", f"user_{idx}")).replace(" ", "_")

            status_text.info(f"Processing {name} ({idx+1}/{total})...")

            temp_ppt = Path(tmpdir) / f"{name}.pptx"
            person_out = base_output / name

            fill_ppt(template_path, temp_ppt, data)
            ppt_to_images(temp_ppt, person_out)

            images = list(person_out.glob("*.PNG"))
            recipient = data.get(email_column)

            if recipient:
                send_email(
                    sender_email,
                    sender_password,
                    recipient,
                    subject=f"{event_name} Certificate",
                    body=f"Dear {data.get('NAME')},\n\nPlease find your certificate.",
                    attachments=images,
                    name=data.get("NAME"),
                    event_name=event_name
                )
                sent_count += 1

            # 🔥 UPDATE PROGRESS
            progress = int(((idx + 1) / total) * 100)
            progress_bar.progress(progress)

    status_text.success(f"✅ {sent_count} / {total} Certificates Sent Successfully!")
    st.success("🎉 All certificates processed!")